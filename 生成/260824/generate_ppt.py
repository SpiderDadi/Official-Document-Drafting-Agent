from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ====== 颜色方案 ======
C_DARK_BLUE = RGBColor(0x1A, 0x3C, 0x5C)
C_MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
C_LIGHT_BLUE = RGBColor(0x4D, 0xA6, 0xFF)
C_ACCENT_GOLD = RGBColor(0xC4, 0x9B, 0x2A)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
C_DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
C_MID_GRAY = RGBColor(0x66, 0x66, 0x66)
C_RED_ACCENT = RGBColor(0xE8, 0x5D, 0x3A)
C_GREEN = RGBColor(0x2D, 0x6A, 0x4F)

def add_bg(slide, color=C_DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
    p.alignment = alignment
    return txBox

def add_bullet_box(slide, left, top, width, height, items, font_size=16, color=C_DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "\u5fae\u8f6f\u96c5\u9ed1"
        p.space_before = Pt(4)
    return txBox


# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, C_DARK_BLUE)
add_shape(slide1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), C_ACCENT_GOLD)
add_shape(slide1, Inches(0.8), Inches(1.0), Inches(0.04), Inches(1.8), C_ACCENT_GOLD)

add_text_box(slide1, Inches(1.2), Inches(0.5), Inches(10), Inches(0.6),
    "\u653f\u7b56\u6df1\u5ea6\u89e3\u8bfb", font_size=16, color=C_ACCENT_GOLD)
add_text_box(slide1, Inches(1.2), Inches(1.2), Inches(11), Inches(1.0),
    "\u53d1\u6539\u59d4\u5bc6\u96c6\u53ec\u5f00\u6c11\u8425\u4f01\u4e1a\u5ea7\u8c08\u4f1a\n\u660e\u786e\u5c06\u51fa\u53f0\u589e\u91cf\u653f\u7b56",
    font_size=36, color=C_WHITE, bold=True)
add_text_box(slide1, Inches(1.2), Inches(2.5), Inches(10), Inches(0.4),
    "\u2014\u2014 \u653f\u7b56\u80cc\u666f\u3001\u5f71\u54cd\u5206\u6790\u4e0e\u673a\u9047\u5c55\u671b", font_size=20, color=C_LIGHT_BLUE)
add_text_box(slide1, Inches(1.2), Inches(3.6), Inches(10), Inches(0.6),
    "2026\u5e748\u670824\u65e5  |  \u653f\u7b56\u7814\u7a76\u8bfe\u9898\u7ec4", font_size=14, color=C_ACCENT_GOLD)
add_shape(slide1, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第2页：目录
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, C_WHITE)
add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
add_text_box(slide2, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
    "\u76ee  \u5f55", font_size=32, color=C_WHITE, bold=True)

toc = [
    ("01", "\u4e8b\u4ef6\u80cc\u666f", "\u4e3a\u4ec0\u4e48\u8fd9\u6b21\u5ea7\u8c08\u4f1a\u503c\u5f97\u9ad8\u5ea6\u5173\u6ce8"),
    ("02", "\u5ea7\u8c08\u4f1a\u6838\u5fc3\u5185\u5bb9", "\u53c2\u4f1a\u4f01\u4e1a\u8bf4\u4e86\u4ec0\u4e48\uff0c\u53d1\u6539\u59d4\u56de\u5e94\u4e86\u4ec0\u4e48"),
    ("03", "\u5b58\u91cf\u4e0e\u589e\u91cf\u653f\u7b56\u5206\u6790", "\u5df2\u51fa\u53f0\u653f\u7b56\u6548\u679c\u5982\u4f55\uff0c\u589e\u91cf\u653f\u7b56\u5c06\u805a\u7126\u4f55\u5904"),
    ("04", "\u5bf9\u6c11\u8425\u4f01\u4e1a\u7684\u5f71\u54cd", "\u878d\u8d44\u3001\u51c6\u5165\u3001\u7ade\u4e89\u3001\u9884\u671f\u56db\u4e2a\u7ef4\u5ea6"),
    ("05", "\u884c\u4e1a\u7ecf\u6d4e\u5f71\u54cd\u8bc4\u4f30", "\u54ea\u4e9b\u884c\u4e1a\u5c06\u7387\u5148\u53d7\u76ca"),
    ("06", "\u653f\u7b56\u5efa\u8bae", "\u4f01\u4e1a\u3001\u6295\u8d44\u8005\u3001\u5730\u65b9\u653f\u5e9c\u5982\u4f55\u5e94\u5bf9"),
]

for i, (num, title, desc) in enumerate(toc):
    y = 1.6 + i * 0.9
    add_text_box(slide2, Inches(1.0), Inches(y), Inches(0.8), Inches(0.5),
        num, font_size=28, color=C_ACCENT_GOLD, bold=True)
    add_text_box(slide2, Inches(2.0), Inches(y), Inches(4), Inches(0.35),
        title, font_size=20, color=C_DARK_BLUE, bold=True)
    add_text_box(slide2, Inches(2.0), Inches(y + 0.35), Inches(8), Inches(0.3),
        desc, font_size=13, color=C_MID_GRAY)
    if i < len(toc) - 1:
        add_shape(slide2, Inches(2.0), Inches(y + 0.75), Inches(9), Inches(0.01), RGBColor(0xE0, 0xE0, 0xE0))

add_shape(slide2, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第3页：事件背景
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, C_WHITE)
add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
add_text_box(slide3, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
    "01  \u4e8b\u4ef6\u80cc\u666f", font_size=32, color=C_WHITE, bold=True)

add_shape(slide3, Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.04), C_ACCENT_GOLD)
add_text_box(slide3, Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.4),
    "\u5173\u952e\u65f6\u95f4\u8282\u70b9", font_size=18, color=C_DARK_BLUE, bold=True)

timeline = [
    "2026\u5e748\u670820\u65e5  |  \u90d1\u6805\u6d01\u4e3b\u4efb\u4e3b\u6301\u53ec\u5f00\u6c11\u4f01\u5ea7\u8c08\u4f1a",
    "5\u5bb6\u6c11\u8425\u4f01\u4e1a\u4ee3\u8868\u53c2\u4f1a\uff08\u7279\u9510\u5fb7\u3001\u695a\u5929\u79d1\u6280\u3001\u8a89\u5b58\u79d1\u6280\u3001\u96f7\u9e1f\u521b\u65b0\u3001\u9102\u5c14\u591a\u65af\u8d44\u6e90\uff09",
    "\u53d1\u6539\u59d4\u660e\u786e\uff1a\u201c\u5145\u5206\u53d1\u6325\u5b58\u91cf\u653f\u7b56\u6548\u80fd\uff0c\u53ca\u65f6\u8c0b\u5212\u51fa\u53f0\u52a1\u5b9e\u7ba1\u7528\u7684\u589e\u91cf\u653f\u7b56\u201d",
    "\u5f3a\u8c03\u534f\u540c\u63a8\u8fdb\u201c\u5341\u4e94\u4e94\u201d\u89c4\u5212\u91cd\u5927\u5de5\u7a0b\u548c\u201c\u516d\u5f20\u7f51\u201d\u89c4\u5212\u5efa\u8bbe",
]
add_bullet_box(slide3, Inches(0.5), Inches(2.3), Inches(5.5), Inches(3.5), timeline, font_size=14, color=C_DARK_TEXT)

add_shape(slide3, Inches(6.5), Inches(1.6), Inches(6.3), Inches(0.04), C_RED_ACCENT)
add_text_box(slide3, Inches(6.5), Inches(1.8), Inches(6.3), Inches(0.4),
    "\u4e3a\u4ec0\u4e48\u91cd\u8981", font_size=18, color=C_DARK_BLUE, bold=True)

add_shape(slide3, Inches(6.5), Inches(2.3), Inches(6.3), Inches(1.5), C_LIGHT_GRAY)
add_bullet_box(slide3, Inches(6.8), Inches(2.4), Inches(5.8), Inches(1.3), [
    "\u2022 \u65f6\u95f4\u70b9\u7279\u6b8a\uff1a\u201c\u5341\u4e94\u4e94\u201d\u5f00\u5c40\u5173\u952e\u7a97\u53e3\u671f",
    "\u2022 \u4fe1\u53f7\u660e\u786e\uff1a\u7a33\u589e\u957f\u653f\u7b56\u52a0\u7801\u4fe1\u53f7\u5f3a\u70c8",
    "\u2022 \u673a\u5236\u5316\uff1a\u5ea7\u8c08\u4f1a\u5e38\u6001\u5316\u91ca\u653e\u5236\u5ea6\u7ea2\u5229",
    "\u2022 \u8986\u76d6\u9762\u5e7f\uff1a5\u5bb6\u4f01\u4e1a\u8986\u76d6\u4e0d\u540c\u884c\u4e1a\u548c\u5730\u533a",
], font_size=14, color=C_DARK_TEXT)

add_shape(slide3, Inches(6.5), Inches(4.0), Inches(6.3), Inches(1.4), RGBColor(0xE8, 0xF4, 0xFD))
add_bullet_box(slide3, Inches(6.8), Inches(4.1), Inches(5.8), Inches(1.2), [
    "\u2022 \u7ecf\u6d4e\u6001\u52bf\uff1a\u201c\u52a8\u80fd\u5411\u65b0\u3001\u7ed3\u6784\u5411\u4f18\u201d",
    "\u2022 \u4f01\u4e1a\u8bc9\u6c42\uff1a\u878d\u8d44\u3001\u534f\u540c\u3001\u58c1\u5792\u3001\u6539\u9020",
    "\u2022 \u653f\u7b56\u7ec4\u5408\uff1a\u201c\u5b58\u91cf\u653f\u7b56+\u589e\u91cf\u653f\u7b56\u201d\u53cc\u8f6e\u9a71\u52a8",
    "\u2022 \u5e02\u573a\u9884\u671f\uff1a\u88ab\u89e3\u8bfb\u4e3a\u65b0\u4e00\u8f6e\u7a33\u589e\u957f\u524d\u594f",
], font_size=14, color=C_DARK_TEXT)

add_shape(slide3, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0), C_DARK_BLUE)
add_text_box(slide3, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.8),
    "\u201c\u4ece\u2018\u4f01\u4e1a\u70b9\u9898\u2019\u5230\u2018\u653f\u5e9c\u7b54\u9898\u2019\uff0c\u5ea7\u8c08\u4f1a\u91ca\u653e\u7684\u662f\u653f\u7b56\u6e29\u5ea6\u7684\u4fe1\u53f7\uff0c\u66f4\u662f\u5236\u5ea6\u521b\u65b0\u7684\u65b9\u5411\u3002\u201d",
    font_size=16, color=C_ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide3, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第4页：座谈会核心内容
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, C_WHITE)
add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
add_text_box(slide4, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
    "02  \u5ea7\u8c08\u4f1a\u6838\u5fc3\u5185\u5bb9", font_size=32, color=C_WHITE, bold=True)

add_shape(slide4, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.04), C_RED_ACCENT)
add_text_box(slide4, Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.4),
    "\u4f01\u4e1a\u8bc9\u6c42\uff085\u4e2a\u65b9\u5411\uff09", font_size=18, color=C_DARK_BLUE, bold=True)

add_bullet_box(slide4, Inches(0.5), Inches(2.3), Inches(5.8), Inches(3.0), [
    "\u2460 \u4f18\u5316\u6295\u8d44\u670d\u52a1 \u2014 \u63d0\u5347\u9879\u76ee\u5ba1\u6279\u6548\u7387",
    "\u2461 \u56fd\u4f01\u6c11\u4f01\u534f\u540c \u2014 \u6253\u7834\u9690\u6027\u58c1\u5792",
    "\u2462 \u89c4\u8303\u7ade\u4e89\u79e9\u5e8f \u2014 \u8425\u9020\u516c\u5e73\u73af\u5883",
    "\u2463 \u62d3\u5bbd\u878d\u8d44\u6e20\u9053 \u2014 \u964d\u4f4e\u878d\u8d44\u6210\u672c",
    "\u2464 \u652f\u6301\u4f20\u7edf\u4ea7\u4e1a\u6539\u9020 \u2014 \u52a0\u5927\u6280\u6539\u652f\u6301",
], font_size=14, color=C_DARK_TEXT)

add_shape(slide4, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.04), C_DARK_BLUE)
add_text_box(slide4, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.4),
    "\u53d1\u6539\u59d4\u8868\u6001\uff084\u4e2a\u65b9\u5411\uff09", font_size=18, color=C_DARK_BLUE, bold=True)

add_bullet_box(slide4, Inches(6.8), Inches(2.3), Inches(6.0), Inches(3.0), [
    "\u25b6 \u5145\u5206\u53d1\u6325\u5b58\u91cf\u653f\u7b56\u6548\u80fd\uff0c\u786e\u4fdd\u7a33\u589e\u957f\u843d\u5730\u89c1\u6548",
    "\u25b6 \u53ca\u65f6\u8c0b\u5212\u51fa\u53f0\u52a1\u5b9e\u7ba1\u7528\u7684\u589e\u91cf\u653f\u7b56",
    "\u25b6 \u534f\u540c\u63a8\u8fdb\u201c\u5341\u4e94\u4e94\u201d\u91cd\u5927\u5de5\u7a0b\u548c\u201c\u516d\u5f20\u7f51\u201d\u5efa\u8bbe",
    "\u25b6 \u6301\u7eed\u4f18\u5316\u8425\u5546\u73af\u5883\uff0c\u89e3\u51b3\u878d\u8d44\u3001\u51c6\u5165\u3001\u7ade\u4e89\u95ee\u9898",
], font_size=14, color=C_DARK_TEXT)

add_shape(slide4, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.3), C_LIGHT_GRAY)
add_bullet_box(slide4, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.1), [
    "\u201c\u4f01\u4e1a\u70b9\u9898 + \u653f\u5e9c\u7b54\u9898\u201d\u6a21\u5f0f\u5f62\u6210\u5236\u5ea6\u95ed\u73af",
    "5\u5bb6\u4f01\u4e1a\u6765\u81ea\u4e0d\u540c\u884c\u4e1a\u548c\u5730\u533a\uff0c\u8bc9\u6c42\u5177\u6709\u5e7f\u6cdb\u4ee3\u8868\u6027",
    "\u53d1\u6539\u59d4\u8868\u6001\u6db5\u76d6\u77ed\u671f\uff08\u5b58\u91cf\uff09+ \u4e2d\u671f\uff08\u589e\u91cf\uff09+ \u957f\u671f\uff08\u5341\u4e94\u4e94\u89c4\u5212\uff09",
], font_size=14, color=C_DARK_TEXT)
add_shape(slide4, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第5页：政策分析
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, C_WHITE)
add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
add_text_box(slide5, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
    "03  \u5b58\u91cf\u4e0e\u589e\u91cf\u653f\u7b56\u5206\u6790", font_size=32, color=C_WHITE, bold=True)

add_shape(slide5, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.5), C_MID_BLUE)
add_text_box(slide5, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.5),
    "\u5b58\u91cf\u653f\u7b56\u6210\u6548", font_size=18, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bullet_box(slide5, Inches(0.5), Inches(2.3), Inches(5.8), Inches(3.0), [
    "\u2022 \u6c11\u8425\u7ecf\u6d4e31\u6761\u63aa\u65bd\u843d\u5730\u6df1\u5316",
    "\u2022 \u51cf\u7a0e\u964d\u8d39\u653f\u7b56\u6301\u7eed\u53d1\u529b",
    "\u2022 \u91d1\u878d\u652f\u6301\u6c11\u8425\u7ecf\u6d4e25\u6761\u9010\u6b65\u89c1\u6548",
    "\u2022 \u8425\u5546\u73af\u5883\u4f18\u5316\u884c\u52a8\u6df1\u5165\u63a8\u8fdb",
    "\u2022 \u4e0a\u534a\u5e74\u6c11\u95f4\u6295\u8d44\u540c\u6bd4+3.2%",
    "\u2022 \u6c11\u4f01\u8fdb\u51fa\u53e3\u603b\u989d\u589e\u957f5.6%",
], font_size=14, color=C_DARK_TEXT)

add_shape(slide5, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5), C_RED_ACCENT)
add_text_box(slide5, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
    "\u589e\u91cf\u653f\u7b56\u65b9\u5411\u9884\u5224", font_size=18, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bullet_box(slide5, Inches(6.8), Inches(2.3), Inches(6.0), Inches(3.0), [
    "\u25b6 \u653e\u5bbd\u6c11\u95f4\u6295\u8d44\u51c6\u5165\uff0c\u805a\u7126\u201c\u516d\u5f20\u7f51\u201d\u5efa\u8bbe",
    "\u25b6 \u52a0\u5927\u7ed3\u6784\u6027\u8d27\u5e01\u653f\u7b56\u5de5\u5177\u652f\u6301\u529b\u5ea6",
    "\u25b6 \u51fa\u53f0\u65b0\u4e00\u8f6e\u51cf\u7a0e\u964d\u8d39\u653f\u7b56",
    "\u25b6 \u5b8c\u5584\u6c11\u4f01\u53c2\u4e0e\u56fd\u5bb6\u91cd\u5927\u6218\u7565\u673a\u5236",
    "\u25b6 \u63a8\u8fdb\u6c11\u8425\u7ecf\u6d4e\u4fc3\u8fdb\u6cd5\u7acb\u6cd5\u8fdb\u7a0b",
    "\u25b6 \u5e38\u6001\u5316\u53ec\u5f00\u6c11\u8425\u4f01\u4e1a\u5ea7\u8c08\u4f1a",
], font_size=14, color=C_DARK_TEXT)

add_shape(slide5, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5), RGBColor(0xE8, 0xF4, 0xFD))
add_text_box(slide5, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.3),
    "\u201c\u516d\u5f20\u7f51\u201d\u89c4\u5212\u5efa\u8bbe \u2014 \u6c11\u4f01\u53c2\u4e0e\u673a\u4f1a\u6e05\u5355",
    font_size=16, color=C_DARK_BLUE, bold=True)
add_text_box(slide5, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5),
    "\u4ea4\u901a\u7f51  |  \u80fd\u6e90\u7f51  |  \u6c34\u5229\u7f51  |  \u4fe1\u606f\u7f51  |  \u7269\u6d41\u7f51  |  \u5e02\u653f\u7f51",
    font_size=18, color=C_MID_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4),
    "\u9884\u8ba1\u91ca\u653e\u6570\u5341\u4e07\u4ebf\u5143\u5e02\u573a\u7a7a\u95f4\uff0c\u6c11\u8425\u4f01\u4e1a\u53ef\u5728\u65b0\u80fd\u6e90\u3001\u6570\u5b57\u7ecf\u6d4e\u3001\u667a\u6167\u57ce\u5e02\u7b49\u9886\u57df\u91cd\u70b9\u53c2\u4e0e",
    font_size=13, color=C_MID_GRAY, alignment=PP_ALIGN.CENTER)
add_shape(slide5, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第6页：对民企影响
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, C_WHITE)
add_shape(slide6, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
add_text_box(slide6, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
    "04  \u5bf9\u6c11\u8425\u4f01\u4e1a\u7684\u5f71\u54cd\u5206\u6790", font_size=32, color=C_WHITE, bold=True)

cards = [
    (Inches(0.5), Inches(1.6), C_MID_BLUE, "\u878d\u8d44\u73af\u5883\u6539\u5584",
     ["\u878d\u8d44\u6e20\u9053\u8fdb\u4e00\u6b65\u62d3\u5bbd", "\u878d\u8d44\u6210\u672c\u7a33\u4e2d\u6709\u964d", "\u79d1\u6280\u578b/\u5236\u9020\u4e1a\u6c11\u4f01\u53d7\u76ca"]),
    (Inches(3.6), Inches(1.6), C_GREEN, "\u5e02\u573a\u51c6\u5165\u5f00\u653e",
     ["\u53c2\u4e0e\u56fd\u5bb6\u91cd\u5927\u9879\u76ee\u5efa\u8bbe", "\u6c11\u95f4\u8d44\u672c\u51c6\u5165\u9886\u57df\u6269\u5927", "\u201c\u516d\u5f20\u7f51\u201d\u91ca\u653e\u65b0\u7a7a\u95f4"]),
    (Inches(6.7), Inches(1.6), C_RED_ACCENT, "\u7ade\u4e89\u79e9\u5e8f\u89c4\u8303",
     ["\u7834\u9664\u9690\u6027\u58c1\u5792", "\u53cd\u5784\u65ad\u7ec6\u5316\u63aa\u65bd", "\u56fd\u4f01\u6c11\u4f01\u516c\u5e73\u7ade\u4e89"]),
    (Inches(9.8), Inches(1.6), C_ACCENT_GOLD, "\u653f\u7b56\u9884\u671f\u7a33\u5b9a",
     ["\u5ea7\u8c08\u4f1a\u5e38\u6001\u5316\u5236\u5ea6\u5316", "\u653f\u7b56\u6c9f\u901a\u673a\u5236\u5b8c\u5584", "\u6295\u8d44\u4fe1\u5fc3\u589e\u5f3a"]),
]

for left, top, color, title, items in cards:
    add_shape(slide6, left, top, Inches(2.8), Inches(0.5), color)
    add_text_box(slide6, left, top, Inches(2.8), Inches(0.5),
        title, font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide6, left, top + Inches(0.5), Inches(2.8), Inches(2.0), C_LIGHT_GRAY)
    add_bullet_box(slide6, left + Inches(0.15), top + Inches(0.6), Inches(2.5), Inches(1.8),
        [f"\u2022 {item}" for item in items], font_size=13, color=C_DARK_TEXT)

add_shape(slide6, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), C_DARK_BLUE)
add_text_box(slide6, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.4),
    "\u6838\u5fc3\u7ed3\u8bba", font_size=18, color=C_ACCENT_GOLD, bold=True)
conclusion_text = (
    "\u6b64\u6b21\u5ea7\u8c08\u4f1a\u91ca\u653e\u7684\u589e\u91cf\u653f\u7b56\u4fe1\u53f7\uff0c\u5c06\u5728\u878d\u8d44\u3001\u51c6\u5165\u3001\u7ade\u4e89\u3001\u9884\u671f\u56db\u4e2a\u7ef4\u5ea6\u5168\u9762\u6539\u5584\u6c11\u8425\u4f01\u4e1a\u8425\u5546\u73af\u5883\u3002\n"
    "\u201c\u5b58\u91cf+\u589e\u91cf\u201d\u53cc\u653f\u7b56\u7ec4\u5408\u62f3\uff0c\u5c06\u5f62\u6210\u4ece\u77ed\u671f\u7ebe\u56f0\u5230\u957f\u671f\u5236\u5ea6\u5efa\u8bbe\u7684\u5b8c\u6574\u653f\u7b56\u94fe\u6761\u3002\n"
    "\u6c11\u8425\u4f01\u4e1a\u5e94\u5bc6\u5207\u5173\u6ce8\u653f\u7b56\u843d\u5730\u8282\u594f\uff0c\u63d0\u524d\u505a\u597d\u6218\u7565\u5e03\u5c40\u548c\u9879\u76ee\u50a8\u5907\u3002"
)
add_text_box(slide6, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.5),
    conclusion_text, font_size=15, color=C_WHITE)
add_shape(slide6, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第7页：行业影响
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, C_WHITE)
add_shape(slide7, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
add_text_box(slide7, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
    "05  \u884c\u4e1a\u7ecf\u6d4e\u5f71\u54cd\u8bc4\u4f30", font_size=32, color=C_WHITE, bold=True)

industries = [
    ("\u65b0\u80fd\u6e90\u4e0e\u7eff\u8272\u4ea7\u4e1a", "\u201c\u80fd\u6e90\u7f51\u201d+\u201c\u5341\u4e94\u4e94\u201d\u80fd\u6e90\u89c4\u5212\u6295\u8d44\u62c9\u52a8", C_GREEN),
    ("\u5148\u8fdb\u5236\u9020\u4e1a\u4e0e\u79d1\u521b", "\u878d\u8d44\u73af\u5883\u6539\u5584+\u6280\u672f\u6539\u9020\u652f\u6301\u653f\u7b56", C_MID_BLUE),
    ("\u57fa\u5efa\u4e0e\u5de5\u7a0b\u673a\u68b0", "\u201c\u516d\u5f20\u7f51\u201d\u91cd\u5927\u5de5\u7a0b\u76f4\u63a5\u62c9\u52a8", C_RED_ACCENT),
    ("\u6570\u5b57\u7ecf\u6d4e\u9886\u57df", "\u201c\u4fe1\u606f\u7f51\u201d+\u65b0\u578b\u57fa\u7840\u8bbe\u65bd\u5efa\u8bbe", C_ACCENT_GOLD),
    ("\u73b0\u4ee3\u7269\u6d41\u4e0e\u5546\u8d38", "\u201c\u7269\u6d41\u7f51\u201d\u89c4\u5212\u5efa\u8bbe\u5e26\u52a8", C_DARK_BLUE),
]

for i, (name, desc, color) in enumerate(industries):
    y = 1.6 + i * 1.05
    add_shape(slide7, Inches(0.5), Inches(y), Inches(0.08), Inches(0.8), color)
    add_text_box(slide7, Inches(0.8), Inches(y), Inches(3.0), Inches(0.4),
        name, font_size=18, color=color, bold=True)
    add_text_box(slide7, Inches(0.8), Inches(y + 0.4), Inches(5.0), Inches(0.4),
        desc, font_size=14, color=C_MID_GRAY)
    stars = "\u2605\u2605\u2605\u2605\u2605" if i < 2 else "\u2605\u2605\u2605\u2605\u2606"
    add_text_box(slide7, Inches(9.5), Inches(y), Inches(3.0), Inches(0.8),
        stars, font_size=18, color=C_ACCENT_GOLD, bold=True)

add_shape(slide7, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), C_LIGHT_GRAY)
add_text_box(slide7, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.7),
    "\u653f\u7b56\u843d\u5730\u8282\u594f\u9884\u5224\uff1a2026\u5e74Q4\u589e\u91cf\u653f\u7b56\u5bc6\u96c6\u51fa\u53f0  \u2192  2027\u5e74\u4e0a\u534a\u5e74\u653f\u7b56\u6548\u5e94\u9010\u6b65\u663e\u73b0  \u2192  \u8d44\u672c\u5e02\u573a\u5df2\u5f62\u6210\u79ef\u6781\u63d0\u632f",
    font_size=14, color=C_DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide7, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第8页：政策建议
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, C_WHITE)
add_shape(slide8, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
add_text_box(slide8, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
    "06  \u653f\u7b56\u5efa\u8bae", font_size=32, color=C_WHITE, bold=True)

suggestions = [
    (Inches(0.5), C_MID_BLUE, "\u5bf9\u4f01\u4e1a\u7684\u5efa\u8bae",
     ["\u5bc6\u5207\u5173\u6ce8\u589e\u91cf\u653f\u7b56\u51fa\u53f0\u8282\u594f", "\u63d0\u524d\u505a\u597d\u9879\u76ee\u50a8\u5907\u548c\u8d44\u91d1\u5b89\u6392",
      "\u805a\u7126\u201c\u516d\u5f20\u7f51\u201d\u673a\u4f1a\u6e05\u5355", "\u52a0\u5f3a\u653f\u4f01\u6c9f\u901a\u5bf9\u63a5"]),
    (Inches(4.7), C_GREEN, "\u5bf9\u6295\u8d44\u8005\u7684\u5efa\u8bae",
     ["\u5173\u6ce8\u65b0\u80fd\u6e90\u3001\u5148\u8fdb\u5236\u9020\u65b9\u5411", "\u628a\u63e1\u6570\u5b57\u7ecf\u6d4e\u653f\u7b56\u7ea2\u5229",
      "\u5173\u6ce8\u6c11\u8425\u7ecf\u6d4e\u4fc3\u8fdb\u6cd5\u7acb\u6cd5", "\u653f\u7b56\u9a71\u52a8\u7684\u6295\u8d44\u673a\u4f1a"]),
    (Inches(8.9), C_RED_ACCENT, "\u5bf9\u5730\u65b9\u653f\u5e9c\u7684\u5efa\u8bae",
     ["\u53ca\u65f6\u4f20\u8fbe\u89e3\u8bfb\u56fd\u5bb6\u653f\u7b56\u7cbe\u795e", "\u5efa\u7acb\u5e38\u6001\u5316\u6c11\u4f01\u5ea7\u8c08\u673a\u5236",
      "\u6e05\u7406\u5404\u7c7b\u9690\u6027\u58c1\u5792", "\u914d\u5408\u201c\u516d\u5f20\u7f51\u201d\u505a\u597d\u9879\u76ee\u50a8\u5907"]),
]

for left, color, title, items in suggestions:
    add_shape(slide8, left, Inches(1.6), Inches(3.9), Inches(0.5), color)
    add_text_box(slide8, left, Inches(1.6), Inches(3.9), Inches(0.5),
        title, font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide8, left, Inches(2.1), Inches(3.9), Inches(3.5), C_LIGHT_GRAY)
    add_bullet_box(slide8, left + Inches(0.2), Inches(2.2), Inches(3.5), Inches(3.3),
        [f"\u2022 {item}" for item in items], font_size=13, color=C_DARK_TEXT)

add_shape(slide8, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), C_DARK_BLUE)
add_text_box(slide8, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8),
    "\u201c\u770b\u61c2\u653f\u7b56\uff0c\u624d\u80fd\u8e29\u51c6\u8282\u62cd\u3002\u201d",
    font_size=24, color=C_ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide8, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ============================================================
# 第9页：尾页
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, C_DARK_BLUE)
add_shape(slide9, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), C_ACCENT_GOLD)

add_text_box(slide9, Inches(1.5), Inches(2.0), Inches(10), Inches(0.6),
    "\u611f\u8c22\u804a\u542c", font_size=40, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide9, Inches(1.5), Inches(2.8), Inches(10), Inches(0.4),
    "\u653f\u7b56\u7814\u7a76\u8bfe\u9898\u7ec4  |  2026\u5e748\u670824\u65e5", font_size=16, color=C_ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
add_text_box(slide9, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
    "\u6570\u636e\u6765\u6e90\uff1a\u56fd\u5bb6\u53d1\u5c55\u548c\u6539\u9769\u59d4\u5458\u4f1a\u5b98\u7f51\uff08ndrc.gov.cn\uff09\n\u58f0\u660e\uff1a\u672c\u62a5\u544a\u4ec5\u4f9b\u53c2\u8003\uff0c\u4e0d\u6784\u6210\u6295\u8d44\u5efa\u8bae",
    font_size=12, color=C_LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
add_shape(slide9, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), C_ACCENT_GOLD)


# ====== 保存 ======
output_path = r"Z:\工作\CC\Official-Document-Drafting-Agent\生成\260824\发改委民企座谈会增量政策解读.pptx"
prs.save(output_path)
print(f"PPT\u5df2\u4fdd\u5b58\u81f3\uff1a{output_path}")