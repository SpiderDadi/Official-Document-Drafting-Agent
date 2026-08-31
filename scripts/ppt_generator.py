#!/usr/bin/env python3
"""
政策 PPT 自动生成器
从政策分析结果自动生成 PPT 文件。

用法：
  python ppt_generator.py --config analysis.json --output 政策解读.pptx

输入格式（与 report_generator.py 相同）：
{
  "date": "2026-08-31",
  "policies": [
    {
      "title": "政策标题",
      "source": "信源",
      "summary": "摘要",
      "rating": "★★★★★",
      "analysis": {
        "summary": "分析总结",
        "analysis": [...],
        "actionItems": [...]
      }
    }
  ]
}

设计原则：
  - 纯代码实现，0 token 消耗
  - 从已有分析结果提取内容，无需 AI 重新理解
  - 统一配色方案（深蓝+金色权威稳重风格）
"""

import json
import argparse
import sys
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("❌ 缺少 python-pptx 库，请运行：pip install python-pptx")
    sys.exit(1)


# ============================================================
# 颜色方案（深蓝+金色，权威稳重风格）
# ============================================================

C_DARK_BLUE = RGBColor(0x1A, 0x3C, 0x5C)
C_MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
C_LIGHT_BLUE = RGBColor(0x4D, 0xA6, 0xFF)
C_ACCENT_GOLD = RGBColor(0xC4, 0x9B, 0x2A)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
C_DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
C_MID_GRAY = RGBColor(0x66, 0x66, 0x66)
C_RED_ACCENT = RGBColor(0xE8, 0x5D, 0x3A)
C_GREEN = RGBColor(0x2D, 0x6A, 0x4F)


# ============================================================
# 辅助函数
# ============================================================

def add_bg(slide, color=C_DARK_BLUE):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, color=C_WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "微软雅黑"
    p.alignment = alignment
    return txBox


def add_bullet_box(slide, left, top, width, height, items,
                   font_size=16, color=C_DARK_TEXT):
    """添加项目符号文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_before = Pt(6)
    return txBox


def add_bottom_bar(slide, color=C_ACCENT_GOLD):
    """添加底部金色装饰条"""
    add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), color)


def add_header(slide, title, section_num=""):
    """添加统一的顶部标题栏"""
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
    if section_num:
        add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                     section_num, font_size=32, color=C_WHITE, bold=True)
    else:
        add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                     title, font_size=32, color=C_WHITE, bold=True)


def safe_text(text, max_len=50):
    """截断超长文本"""
    if not text:
        return ""
    text = text.replace("\n", " ").replace("\r", "")
    return text[:max_len] + ("..." if len(text) > max_len else "")


# ============================================================
# PPT 生成主函数
# ============================================================

def generate_ppt(config):
    """根据政策分析结果生成 PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    date_str = config.get("date", datetime.now().strftime("%Y-%m-%d"))
    policies = config.get("policies", [])

    # ============================================================
    # 第 1 页：封面
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_DARK_BLUE)
    add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), C_ACCENT_GOLD)
    add_shape(slide, Inches(0.8), Inches(1.0), Inches(0.04), Inches(1.8), C_ACCENT_GOLD)

    add_text_box(slide, Inches(1.2), Inches(0.5), Inches(10), Inches(0.6),
                 "政策深度解读", font_size=16, color=C_ACCENT_GOLD)

    main_title = f"{date_str} 重点政策分析"
    if policies:
        main_title = safe_text(policies[0].get("title", ""), 40)
    add_text_box(slide, Inches(1.2), Inches(1.2), Inches(11), Inches(1.0),
                 main_title, font_size=36, color=C_WHITE, bold=True)

    subtitle = f"共 {len(policies)} 条重点政策"
    add_text_box(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(0.4),
                 f"—— {subtitle}", font_size=20, color=C_LIGHT_BLUE)

    add_text_box(slide, Inches(1.2), Inches(3.6), Inches(10), Inches(0.6),
                 f"{date_str}  |  政策研究课题组", font_size=14, color=C_ACCENT_GOLD)

    add_bottom_bar(slide)

    # ============================================================
    # 第 2 页：目录（只生成前3条）
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), C_DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                 "目  录", font_size=32, color=C_WHITE, bold=True)

    toc_items = []
    for i, policy in enumerate(policies[:5]):
        num = f"{i+1:02d}"
        title = safe_text(policy.get("title", ""), 30)
        rating = policy.get("rating", "")
        desc = f"{policy.get('source', '')}  {rating}"
        toc_items.append((num, title, desc))

    for i, (num, title, desc) in enumerate(toc_items):
        y = 1.6 + i * 0.9
        add_text_box(slide, Inches(1.0), Inches(y), Inches(0.8), Inches(0.5),
                     num, font_size=28, color=C_ACCENT_GOLD, bold=True)
        add_text_box(slide, Inches(2.0), Inches(y), Inches(4), Inches(0.35),
                     title, font_size=20, color=C_DARK_BLUE, bold=True)
        add_text_box(slide, Inches(2.0), Inches(y + 0.35), Inches(8), Inches(0.3),
                     desc, font_size=13, color=C_MID_GRAY)
        if i < len(toc_items) - 1:
            add_shape(slide, Inches(2.0), Inches(y + 0.75), Inches(9), Inches(0.01),
                      RGBColor(0xE0, 0xE0, 0xE0))

    add_bottom_bar(slide)

    # ============================================================
    # 为每条政策生成 2-3 页内容
    # ============================================================
    for i, policy in enumerate(policies[:3]):  # 最多为前3条政策生成详细页
        title = safe_text(policy.get("title", ""), 30)
        source = policy.get("source", "")
        rating = policy.get("rating", "")
        section_num = f"{i+1:02d}"
        analysis = policy.get("analysis", {})

        if isinstance(analysis, str):
            # 如果是字符串，简单处理
            analysis = {"summary": analysis, "qa": [], "actionItems": []}
        elif not isinstance(analysis, dict):
            analysis = {"summary": "", "qa": [], "actionItems": []}

        qa_list = analysis.get("qa", analysis.get("analysis", []))
        action_items = analysis.get("actionItems", [])

        # --- 政策详情页 A：核心要点 ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, C_WHITE)
        add_header(slide, title, f"{section_num}  {source}  {rating}")

        # 左侧：摘要
        add_shape(slide, Inches(0.5), Inches(1.6), Inches(0.04), Inches(5.5), C_ACCENT_GOLD)

        summary = safe_text(policy.get("summary", "") or analysis.get("summary", ""), 100)
        add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4),
                     "政策摘要", font_size=18, color=C_DARK_BLUE, bold=True)

        summary_items = [s.strip() for s in summary.replace("。", "。\n").replace("；", "；\n").split("\n") if s.strip()]
        add_bullet_box(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5),
                       summary_items[:6], font_size=14, color=C_DARK_TEXT)

        # 右侧：关键要点
        key_points = safe_text(policy.get("keyPoints", ""), 80)
        add_shape(slide, Inches(6.5), Inches(1.6), Inches(6.3), Inches(0.04), C_RED_ACCENT)
        add_text_box(slide, Inches(6.8), Inches(1.8), Inches(6.3), Inches(0.4),
                     "关键要点", font_size=18, color=C_DARK_BLUE, bold=True)

        if key_points:
            kp_items = [s.strip() for s in key_points.replace("。", "。\n").replace("；", "；\n").split("\n") if s.strip()]
            add_bullet_box(slide, Inches(6.8), Inches(2.3), Inches(6.3), Inches(3.5),
                           kp_items[:5], font_size=14, color=C_DARK_TEXT)

        # 底部金句
        if analysis.get("summary"):
            quote = safe_text(analysis["summary"], 80)
            add_shape(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0), C_DARK_BLUE)
            add_text_box(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.8),
                         f""{quote}"", font_size=16, color=C_ACCENT_GOLD, bold=True,
                         alignment=PP_ALIGN.CENTER)

        add_bottom_bar(slide)

        # --- 政策详情页 B：深度分析 ---
        if qa_list:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide, C_WHITE)
            add_header(slide, f"{title} — 深度分析", f"{section_num}.B  三轮递进分析")

            # 按轮次分组
            rounds = {}
            for qa in qa_list:
                r = qa.get("round", "综合")
                if r not in rounds:
                    rounds[r] = []
                rounds[r].append(qa)

            # 分两栏展示
            round_list = list(rounds.items())

            # 左栏
            if len(round_list) > 0:
                round_name, round_qas = round_list[0]
                add_shape(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.04), C_ACCENT_GOLD)
                add_text_box(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.4),
                             round_name, font_size=16, color=C_DARK_BLUE, bold=True)

                left_items = []
                for qa in round_qas[:3]:
                    q = safe_text(qa.get("question", ""), 40)
                    a = safe_text(qa.get("answer", ""), 60)
                    if q and a:
                        left_items.append(f"{q}：{a}")
                add_bullet_box(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(3.5),
                               left_items, font_size=13, color=C_DARK_TEXT)

            # 右栏
            if len(round_list) > 1:
                round_name, round_qas = round_list[1]
                add_shape(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.04), C_MID_BLUE)
                add_text_box(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.4),
                             round_name, font_size=16, color=C_DARK_BLUE, bold=True)

                right_items = []
                for qa in round_qas[:3]:
                    q = safe_text(qa.get("question", ""), 40)
                    a = safe_text(qa.get("answer", ""), 60)
                    if q and a:
                        right_items.append(f"{q}：{a}")
                add_bullet_box(slide, Inches(6.8), Inches(2.3), Inches(6.0), Inches(3.5),
                               right_items, font_size=13, color=C_DARK_TEXT)

            # 第三轮（如果有）
            if len(round_list) > 2:
                round_name, round_qas = round_list[2]
                add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.04), C_RED_ACCENT)
                add_text_box(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.3),
                             round_name, font_size=16, color=C_DARK_BLUE, bold=True)

                third_items = []
                for qa in round_qas[:3]:
                    q = safe_text(qa.get("question", ""), 40)
                    a = safe_text(qa.get("answer", ""), 60)
                    if q and a:
                        third_items.append(f"{q}：{a}")
                add_bullet_box(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.0),
                               third_items, font_size=13, color=C_DARK_TEXT)

            add_bottom_bar(slide)

        # --- 政策详情页 C：行动建议 ---
        if action_items:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide, C_WHITE)
            add_header(slide, f"{title} — 行动建议", f"{section_num}.C  行动建议")

            # 最多显示6条
            cards_per_row = min(3, len(action_items))
            card_width = 3.9
            spacing = 0.3

            for j, action in enumerate(action_items[:6]):
                row = j // cards_per_row
                col = j % cards_per_row
                x = 0.5 + col * (card_width + spacing)
                y = 1.6 + row * 2.5

                colors = [C_MID_BLUE, C_GREEN, C_RED_ACCENT, C_ACCENT_GOLD, C_DARK_BLUE, C_LIGHT_BLUE]
                card_color = colors[j % len(colors)]

                add_shape(slide, Inches(x), Inches(y), Inches(card_width), Inches(0.5), card_color)
                add_text_box(slide, Inches(x), Inches(y), Inches(card_width), Inches(0.5),
                             f"建议 {j+1}", font_size=16, color=C_WHITE, bold=True,
                             alignment=PP_ALIGN.CENTER)
                add_shape(slide, Inches(x), Inches(y + 0.5), Inches(card_width), Inches(1.8), C_LIGHT_GRAY)
                add_text_box(slide, Inches(x + 0.2), Inches(y + 0.6), Inches(card_width - 0.4), Inches(1.6),
                             safe_text(action, 100), font_size=14, color=C_DARK_TEXT)

            add_bottom_bar(slide)

    # ============================================================
    # 最后一页：尾页
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_DARK_BLUE)
    add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), C_ACCENT_GOLD)

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.6),
                 "感谢聆听", font_size=40, color=C_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.4),
                 f"政策研究课题组  |  {date_str}", font_size=16, color=C_ACCENT_GOLD,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
                 f"数据来源：{', '.join(set(p.get('source', '') for p in policies if p.get('source')))}\n"
                 f"本报告由 AI 深度分析 + 自动生成系统产出，仅供参考",
                 font_size=12, color=C_LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide)

    return prs


def main():
    parser = argparse.ArgumentParser(description="政策 PPT 自动生成器")
    parser.add_argument('--config', '-c', type=str,
                        help='JSON 配置文件路径')
    parser.add_argument('--json', '-j', type=str,
                        help='JSON 配置字符串')
    parser.add_argument('--output', '-o', type=str, required=True,
                        help='输出 PPT 文件路径')

    args = parser.parse_args()

    if args.config:
        config = json.loads(Path(args.config).read_text(encoding="utf-8"))
    elif args.json:
        config = json.loads(args.json)
    else:
        print("错误：请提供 --config 或 --json 参数", file=sys.stderr)
        sys.exit(1)

    # 如果没有 date，自动填充
    if "date" not in config:
        config["date"] = datetime.now().strftime("%Y-%m-%d")

    # 如果没有 policies，尝试转换
    if "policies" not in config:
        if "title" in config:
            config["policies"] = [config]
        elif "analyses" in config:
            config["policies"] = config["analyses"]

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = generate_ppt(config)
    prs.save(str(output_path))
    print(f"✅ PPT 已保存：{output_path}")

    policies = config.get("policies", [])
    slide_count = len(prs.slides)
    print(f"   共 {slide_count} 页，分析 {len(policies[:3])} 条政策")


if __name__ == "__main__":
    main()
